"""Generate slides.pptx (deliverable). Minimal, replace with the talkshow/editorial deck."""
from pathlib import Path
HERE = Path(__file__).resolve().parent

def build(metrics, segments):
    from pptx import Presentation
    prs = Presentation()
    def slide(title, body):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        s.placeholders[1].text = body
    slide("Churn Prediction & Segmentation", "Digital Bank — Business Analytics group project")
    slide("Problem & KPIs", "Who churns next quarter, in which segment?\nKPIs: churn rate, CLV, retention cost")
    slide("Model results", f"AUC-ROC: {metrics['churn'].get('auc_roc')}\n"
          f"Recall (churn): {metrics['churn'].get('recall')}\nF1: {metrics['churn'].get('f1')}")
    personas = segments.get("personas", {})
    by_seg = segments.get("churn_rate_by_segment", {})
    slide("Segments & churn", "\n".join(f"{k}: {personas.get(str(k),'?')} — churn {v}"
                                        for k, v in by_seg.items()) or "TBD")
    slide("Retention strategy", "Per-segment actions tied to CLV/cost (see report §5)")
    slide("Limitations & next steps", "Synthetic data realism; threshold tuning; live retraining")
    out = HERE / "slides.pptx"; prs.save(out)
    print(f"[M8] slides -> {out}")

if __name__ == "__main__":
    import json
    build(json.load(open(HERE.parent/"research/outputs/metrics/metrics.json")),
          json.load(open(HERE.parent/"research/outputs/metrics/segments.json")))
